import streamlit as st
from streamlit_option_menu import option_menu
import ollama
from docx import Document
from docx.shared import Pt
from docx.enum.style import WD_STYLE_TYPE
from io import BytesIO
import markdown
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt as PptxPt

# Set page config for a wider layout
st.set_page_config(layout="wide", page_title="Educational Assistant")

def generate_response(prompt, context=None):
    if context:
        messages = [{'role': 'user', 'content': f"Based on the following content, respond to this prompt: {prompt}\n\nContent: {context}"}]
    else:
        messages = [{'role': 'user', 'content': prompt}]
    
    response = ollama.chat(model='gemma2:2b', messages=messages)
    return response['message']['content']



def export_to_word(content):
    # Convert markdown to HTML
    html = markdown.markdown(content)
    
    # Parse HTML
    soup = BeautifulSoup(html, 'html.parser')
    
    # Create a new Document
    doc = Document()
    
    # Define styles
    styles = doc.styles
    style_normal = styles['Normal']
    style_heading = styles.add_style('Heading', WD_STYLE_TYPE.PARAGRAPH)
    style_bold = styles.add_style('Bold', WD_STYLE_TYPE.CHARACTER)
    style_italic = styles.add_style('Italic', WD_STYLE_TYPE.CHARACTER)
    
    style_heading.font.size = Pt(16)
    style_heading.font.bold = True
    style_bold.font.bold = True
    style_italic.font.italic = True
    
    # Function to add formatted text
    def add_formatted_text(element, paragraph):
        for child in element.children:
            if child.name == 'strong':
                paragraph.add_run(child.text).bold = True
            elif child.name == 'em':
                paragraph.add_run(child.text).italic = True
            elif child.name is None:  # This is a text node
                paragraph.add_run(child.string)
    
    # Process HTML elements
    for element in soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'ul', 'ol']):
        if element.name.startswith('h'):
            paragraph = doc.add_paragraph(style=style_heading)
            add_formatted_text(element, paragraph)
        elif element.name == 'p':
            paragraph = doc.add_paragraph(style=style_normal)
            add_formatted_text(element, paragraph)
        elif element.name in ['ul', 'ol']:
            for i, li in enumerate(element.find_all('li'), start=1):
                paragraph = doc.add_paragraph(style=style_normal)
                paragraph.paragraph_format.left_indent = Pt(20)
                run = paragraph.add_run('• ' if element.name == 'ul' else f"{i}. ")
                add_formatted_text(li, paragraph)
    
    bio = BytesIO()
    doc.save(bio)
    return bio.getvalue()

def create_powerpoint(content):
    # Generate slide content using Gemma
    slide_content = generate_response(f"Create a 5-slide PowerPoint presentation based on this content. For each slide, provide a title prefixed with 'title:' and 3-5 bullet points prefixed with 'bullet points:'. Use markdown formatting. Separate each slide with '---': {content}")
    
    # Create a new PowerPoint presentation
    prs = Presentation()
    
    # Parse the generated content
    slides = slide_content.split('---')
    
    for slide_content in slides:
        if not slide_content.strip():
            continue  # Skip empty slides
        
        # Split the content into title and bullet points
        parts = slide_content.split('bullet points:', 1)
        if len(parts) != 2:
            continue  # Skip slides with incorrect format
        
        title_part = parts[0].strip()
        bullet_points_part = parts[1].strip()
        
        # Extract title and convert from markdown
        title_html = markdown.markdown(title_part.replace('title:', '').strip())
        title_soup = BeautifulSoup(title_html, 'html.parser')
        title = title_soup.get_text().strip()
        
        # Convert bullet points from markdown to HTML, then to plain text
        bullet_points_html = markdown.markdown(bullet_points_part)
        bullet_soup = BeautifulSoup(bullet_points_html, 'html.parser')
        bullet_points = [li.get_text().strip() for li in bullet_soup.find_all('li')]
        
        # Add a slide
        slide_layout = prs.slide_layouts[1]  # Using the bullet slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set the title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Add bullet points
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        
        if bullet_points:
            tf.text = bullet_points[0]
            
            for point in bullet_points[1:]:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
        else:
            tf.text = "No bullet points provided."
    
    # Ensure at least one slide is created
    if not prs.slides:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = "Presentation"
        subtitle = slide.placeholders[1]
        subtitle.text = "No content was generated for this presentation."
    
    # Save the presentation
    pptx_stream = BytesIO()
    prs.save(pptx_stream)
    return pptx_stream.getvalue()

# Custom CSS for better aesthetics with improved readability
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Poppins:wght@400;600&display=swap');

    :root {
        --primary-color: #6a11cb;
        --secondary-color: #2575fc;
        --accent-color: #ffd166;
        --text-color: #ffffff;
        --background-color: rgba(255, 255, 255, 0.1);
    }

    .stApp {
        background: linear-gradient(135deg, var(--primary-color), var(--secondary-color));
        font-family: 'Poppins', sans-serif;
    }

    .main {
        max-width: 800px;
        margin: 0 auto;
        padding: 2rem;
        background-color: var(--background-color);
        backdrop-filter: blur(10px);
        border-radius: 10px;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
    }

    /* Sidebar styling */
    [data-testid="stSidebar"] {
        background: linear-gradient(135deg, var(--primary-color), var(--secondary-color));
        padding: 2rem 1rem;
    }

    [data-testid="stSidebar"] .stSelectbox {
        margin-bottom: 1rem;
    }

    /* Style for the menu items in sidebar */
    .css-1544g2n {
        padding: 1rem;
        background-color: var(--background-color);
        border-radius: 10px;
        margin-bottom: 1rem;
    }

    /* Chat prompt box styling */
    .stTextInput > div > div > input {
        background-color: var(--background-color);
        color: var(--text-color);
        border: 1px solid rgba(255, 255, 255, 0.2);
        border-radius: 5px;
        padding: 0.5rem 1rem;
    }

    /* File uploader styling */
    .stFileUploader {
        background-color: var(--background-color);
        border: 1px solid rgba(255, 255, 255, 0.2);
        border-radius: 10px;
        padding: 1rem;
    }

    h1, h2, h3 {
        color: var(--accent-color);
        font-weight: 600;
        margin-bottom: 1rem;
    }

    p {
        color: var(--text-color);
        line-height: 1.6;
        margin-bottom: 1rem;
    }

    .stButton > button {
        background-color: var(--accent-color);
        color: var(--primary-color);
        font-weight: 600;
        border: none;
        border-radius: 5px;
        padding: 0.5rem 1rem;
        transition: all 0.3s ease;
    }

    .stButton > button:hover {
        background-color: var(--primary-color);
        color: var(--accent-color);
    }

    .chat-message {
        padding: 1.5rem;
        border-radius: 0.5rem;
        margin-bottom: 1rem;
        display: flex;
        box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
        background-color: var(--background-color);
    }

    .chat-message .message {
       width: 100%;
       color: var(--text-color);
    }

    /* Custom scrollbar */
    ::-webkit-scrollbar {
        width: 10px;
    }

    ::-webkit-scrollbar-track {
        background: rgba(255, 255, 255, 0.1);
    }

    ::-webkit-scrollbar-thumb {
        background: var(--accent-color);
        border-radius: 5px;
    }

    ::-webkit-scrollbar-thumb:hover {
        background: var(--primary-color);
    }
</style>
""", unsafe_allow_html=True)

st.title("✨ AIED Special Sauce Mega Platform V2.0 ✨")

# Sidebar
with st.sidebar:
    selected = option_menu("Menu", ["Chat", "Actions", "Export"], 
        icons=['chat', 'gear', 'file-earmark-arrow-down'], menu_icon="cast", default_index=0,
        styles={
            "container": {"padding": "1rem", "background-color": "rgba(255, 255, 255, 0.1)"},
            "icon": {"color": "var(--accent-color)", "font-size": "25px"}, 
            "nav-link": {"color": "var(--text-color)", "font-size": "16px", "text-align": "left", "margin":"0px", "--hover-color": "rgba(255, 255, 255, 0.2)"},
            "nav-link-selected": {"background-color": "var(--accent-color)", "color": "var(--primary-color)"},
        }
    )

# Initialize chat history
if "messages" not in st.session_state:
    st.session_state.messages = []

if selected == "Chat":
    # File uploader for Q&A
    uploaded_file = st.file_uploader("📄 Upload a document for context (optional)", type=["txt"])
    if uploaded_file:
        file_content = uploaded_file.getvalue().decode("utf-8")
        st.session_state.file_content = file_content
        st.success("✅ File uploaded successfully!")

    # Display chat messages from history on app rerun
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

    # Accept user input
    if prompt := st.chat_input("💬 What would you like to do today?"):
        # Add user message to chat history
        st.session_state.messages.append({"role": "user", "content": prompt})
        # Display user message in chat message container
        with st.chat_message("user"):
            st.markdown(prompt)

        # Generate response
        context = st.session_state.file_content if "file_content" in st.session_state else None
        response = generate_response(prompt, context)

        # Display assistant response in chat message container
        with st.chat_message("assistant"):
            st.markdown(response)
        
        # Add assistant response to chat history
        st.session_state.messages.append({"role": "assistant", "content": response})

        # Force a rerun to update the display immediately
        st.experimental_rerun()

elif selected == "Actions":
    if st.session_state.messages:
        latest_message = st.session_state.messages[-1]["content"]

        action = st.selectbox("🔧 Choose an action:", [
            "Select an action",
            "Translate to French",
            "Create a lesson plan",
            "Generate a vocabulary list",
            "Add sparkle to the text",
            "Create a PowerPoint presentation"
        ])

        if st.button("▶️ Perform Action"):
            if action == "Translate to French":
                prompt = f"Translate the following text to French: {latest_message}"
            elif action == "Create a lesson plan":
                prompt = f"Produce a clear, single 50 minute lesson plan based on the following content. Provide 3 learning objectives and success criteria: {latest_message}"
            elif action == "Generate a vocabulary list":
                prompt = f"Generate a vocabulary list based on the following content: {latest_message}"
            elif action == "Add sparkle to the text":
                prompt = f"Take this, keep the text the same, but cover it in relevant emojis: {latest_message}"
            elif action == "Create a PowerPoint presentation":
                pptx_data = create_powerpoint(latest_message)
                st.download_button(
                    label="📥 Download PowerPoint",
                    data=pptx_data,
                    file_name="presentation.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
                st.success("✅ PowerPoint presentation created successfully!")
                st.stop()
            else:
                st.warning("⚠️ Please select an action.")
                st.stop()

            response = generate_response(prompt)
            
            # Display the result
            with st.chat_message("assistant"):
                st.markdown(response)
            
            # Add the result to chat history
            st.session_state.messages.append({"role": "assistant", "content": response})

    else:
        st.warning("⚠️ No chat history. Please start a conversation first.")

elif selected == "Export":
    if st.session_state.messages:
        latest_message = st.session_state.messages[-1]["content"]
        
        export_type = st.radio("📁 Choose export format:", ["Word", "PowerPoint"])
        
        if export_type == "Word":
            if st.button("📄 Export to Word"):
                word_doc = export_to_word(latest_message)
                st.download_button(
                    label="📥 Download Word Document",
                    data=word_doc,
                    file_name="exported_content.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                )
        else:  # PowerPoint
            if st.button("📊 Export to PowerPoint"):
                pptx_data = create_powerpoint(latest_message)
                st.download_button(
                    label="📥 Download PowerPoint",
                    data=pptx_data,
                    file_name="presentation.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
    else:
        st.warning("⚠️ No content to export. Please generate some content first.")
